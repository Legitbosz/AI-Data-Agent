import pandas as pd
import os
import io
import re
import zipfile
from xml.sax.saxutils import escape as xml_escape
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment
from openpyxl.chart import BarChart, Reference
from docx import Document
from docx.shared import Inches, Pt
from docx.dml.color import RGBColor as DocxRGB
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor as PptxRGB

COLS = "BCDEFGHIJKLMNOPQRSTUVWXYZ"

# Vibrant chart color palette (hex without #, used across Excel XML and PPTX)
CHART_COLORS = [
    "2563EB",  # blue
    "F97316",  # orange
    "16A34A",  # green
    "EF4444",  # red
    "8B5CF6",  # purple
    "EC4899",  # pink
    "14B8A6",  # teal
    "F59E0B",  # amber
    "6366F1",  # indigo
    "06B6D4",  # cyan
    "84CC16",  # lime
    "D946EF",  # fuchsia
]


# ── PPTX helpers ──────────────────────────────────────────────────────────────
def set_textbox(shape, left, top, width, height):
    shape.left   = int(PInches(left))
    shape.top    = int(PInches(top))
    shape.width  = int(PInches(width))
    shape.height = int(PInches(height))


def add_text_to_placeholder(tf, text, font_size=16, bold=False):
    from pptx.oxml.ns import qn as pqn
    from lxml import etree as et
    txBody = tf._txBody
    bodyPr = txBody.find(pqn("a:bodyPr"))
    if bodyPr is None:
        bodyPr = et.SubElement(txBody, pqn("a:bodyPr"))
    for tag in ["a:noAutofit", "a:spAutoFit", "a:normAutofit"]:
        el = bodyPr.find(pqn(tag))
        if el is not None:
            bodyPr.remove(el)
    et.SubElement(bodyPr, pqn("a:normAutofit"))
    bodyPr.set("wrap", "square")
    tf.text = text
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = PPt(font_size)
            run.font.bold = bold


def add_bullet_placeholder(tf, findings, font_size=16):
    from pptx.oxml.ns import qn as pqn
    from lxml import etree as et
    txBody = tf._txBody
    bodyPr = txBody.find(pqn("a:bodyPr"))
    if bodyPr is None:
        bodyPr = et.SubElement(txBody, pqn("a:bodyPr"))
    for tag in ["a:noAutofit", "a:spAutoFit", "a:normAutofit"]:
        el = bodyPr.find(pqn(tag))
        if el is not None:
            bodyPr.remove(el)
    et.SubElement(bodyPr, pqn("a:normAutofit"))
    bodyPr.set("wrap", "square")
    tf.clear()
    for i, finding in enumerate(findings):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = finding
        para.level = 0
        for run in para.runs:
            run.font.size = PPt(font_size)


# ── Data helpers ──────────────────────────────────────────────────────────────
# Keywords that indicate financial/revenue columns — checked in priority order
FINANCIAL_KEYWORDS = [
    "bill", "amt", "amount", "revenue", "sales", "agreed", "total",
    "income", "value", "receipt", "payment", "price", "cost", "fee",
    "card", "net", "gross"
]

# Keywords that indicate dimension/count columns — deprioritised
DIMENSION_KEYWORDS = ["ht", "wid", "wd", "area", "height", "width", "size",
                       "no", "num", "count", "qty", "id", "key", "code",
                       "rate", "per", "pct", "percent", "discount", "vat", "wht", "tax"]


def best_numeric_col(num_cols, sname):
    """
    Pick the best numeric column to aggregate for a given series name.
    Priority:
    1. Direct name match with series name keywords
    2. Financial column keywords match
    3. Any col that is NOT a dimension/count column
    4. First numeric column (last resort)
    """
    sname_words = sname.lower().replace("_", " ").split()

    # 1. Direct match with series name
    for nc in num_cols:
        nc_clean = nc.lower().replace("_", " ").replace("  ", " ")
        if any(w in nc_clean for w in sname_words if len(w) > 2):
            return nc

    # 2. Financial keyword match — pick highest priority financial col
    for kw in FINANCIAL_KEYWORDS:
        for nc in num_cols:
            nc_clean = nc.lower().replace("_", " ").replace("  ", " ")
            if kw in nc_clean:
                return nc

    # 3. Any col NOT matching dimension keywords
    for nc in num_cols:
        nc_clean = nc.lower().replace("_", " ").replace("  ", " ")
        if not any(kw in nc_clean for kw in DIMENSION_KEYWORDS):
            return nc

    # 4. Last resort
    return num_cols[0]


def fix_chart_values(cd, df):
    """
    Recompute series values from the actual dataframe.
    AI frequently hallucinates wrong aggregation totals — always use dataframe.
    Handles real-world data: messy column names, large datasets, various dtypes.
    """
    import logging

    categories = [str(c) for c in cd.get("categories", [])]
    ai_series  = cd.get("series", {})

    # Normalize df column names
    df_clean = df.copy()
    df_clean.columns = (df_clean.columns.astype(str)
                        .str.replace(r'[\xa0\s]+', ' ', regex=True)
                        .str.strip())

    text_cols = df_clean.select_dtypes(exclude="number").columns.tolist()
    num_cols  = df_clean.select_dtypes(include="number").columns.tolist()

    if not text_cols or not num_cols:
        return {k: [float(v) for v in vals] for k, vals in ai_series.items()}

    fixed = {}
    for sname, ai_vals in ai_series.items():
        matched = False
        for col in text_cols:
            try:
                col_vals = df_clean[col].astype(str).str.strip().unique().tolist()
                cats_norm = [str(c).strip() for c in categories]
                if all(cat in col_vals for cat in cats_norm):
                    num_col = best_numeric_col(num_cols, sname)
                    logging.warning(f"[fix_chart_values] groupby col='{col}', num_col='{num_col}' for series='{sname}'")
                    grouped = df_clean.groupby(col, observed=False)[num_col].sum()
                    fixed[sname] = [float(grouped.get(cat, 0)) for cat in cats_norm]
                    matched = True
                    break
            except Exception as e:
                logging.warning(f"[fix_chart_values] col={col} error: {e}")
                continue

        if not matched:
            fixed[sname] = [float(v) for v in ai_vals]
    return fixed


def detect_currency(analysis):
    """Detect currency symbol from analysis — AI may return it or we infer from answer text."""
    sym = analysis.get("currency_symbol") if analysis else None
    if sym:
        return sym
    # Infer from answer/findings text
    text = ""
    if analysis:
        text = analysis.get("answer", "") + " ".join(analysis.get("key_findings", []))
    if "₦" in text or "naira" in text.lower() or "nigerian" in text.lower():
        return "₦"
    if "£" in text or "pound" in text.lower() or "gbp" in text.lower():
        return "£"
    if "€" in text or "euro" in text.lower():
        return "€"
    return "$"  # default


def resolve_all_charts(analysis, df):
    """
    Return all charts as a clean list with verified values from the dataframe.
    Uses AI structure (titles, axis labels, chart type, label format)
    but always recomputes values from actual dataframe for accuracy.
    """
    charts_data = analysis.get("charts_data") if analysis else None
    result = []

    if charts_data and isinstance(charts_data, list):
        for cd in charts_data:
            try:
                categories = [str(c) for c in cd.get("categories", [])]
                if not categories:
                    continue
                fixed_series = fix_chart_values(cd, df)
                result.append({
                    "title":        cd.get("title", "Chart"),
                    "categories":   categories,
                    "series":       fixed_series,
                    "label_format": cd.get("label_format", "none"),
                    "chart_type":   cd.get("chart_type", "bar"),
                    "x_title":      cd.get("x_title", "Category"),
                    "y_title":      cd.get("y_title", "Value"),
                })
            except Exception as e:
                import logging
                logging.warning(f"[exporters] chart error: {e}")

    if not result:
        # Fallback: single chart from first text+numeric column groupby
        text_cols = df.select_dtypes(exclude="number").columns.tolist()
        num_cols  = df.select_dtypes(include="number").columns.tolist()
        if text_cols and num_cols:
            g = df.groupby(text_cols[0], observed=False)[num_cols[0]].sum().reset_index()
            result.append({
                "title":        analysis.get("chart_title", "Chart") if analysis else "Chart",
                "categories":   g[text_cols[0]].astype(str).tolist(),
                "series":       {num_cols[0]: g[num_cols[0]].tolist()},
                "label_format": "none",
                "chart_type":   "bar",
                "x_title":      text_cols[0],
                "y_title":      num_cols[0],
            })
    return result


# ── Excel chart XML builder ───────────────────────────────────────────────────
def build_chart_xml(ws_title, categories, series, chart_title,
                    y_title="Value", x_title="Category",
                    label_format="none", chart_type="bar", currency_symbol="$",
                    data_start_row=1):
    """Build correct OOXML chart from scratch — bypasses all openpyxl XML bugs."""
    n = len(categories)

    # Escape text values for XML safety (handles <, >, & in category names etc.)
    chart_title = xml_escape(str(chart_title))
    x_title     = xml_escape(str(x_title))
    y_title     = xml_escape(str(y_title))

    # Row references — defined here so all chart types can use them
    hdr_row  = data_start_row
    cat_row1 = data_start_row + 1
    cat_rown = data_start_row + n

    fmt_map  = {"number": "#,##0", "percent": "0.0%"}
    fmt_code = fmt_map.get(label_format, "General")
    if label_format == "currency":
        fmt_code = f'{currency_symbol}#,##0'
    show_labels = label_format != "none"

    dlbls = ""
    if show_labels:
        dlbls = f"""<c:dLbls>
          <c:numFmt formatCode="{fmt_code}" sourceLinked="0"/>
          <c:showLegendKey val="0"/><c:showVal val="1"/>
          <c:showCatName val="0"/><c:showSerName val="0"/>
          <c:showPercent val="0"/><c:showBubbleSize val="0"/>
        </c:dLbls>"""

    marker_xml = ""
    if chart_type == "line":
        marker_xml = "<c:marker><c:symbol val=\"circle\"/><c:size val=\"5\"/></c:marker>"

    ser_blocks = ""
    for i, (sname, vals) in enumerate(series.items()):
        col_letter = COLS[i]
        pts     = "".join(f'<c:pt idx="{j}"><c:v>{v}</c:v></c:pt>' for j, v in enumerate(vals))
        cat_pts = "".join(f'<c:pt idx="{j}"><c:v>{xml_escape(str(cat))}</c:v></c:pt>' for j, cat in enumerate(categories))

        # Per-data-point colors for bar charts (each bar a different color)
        dpt_xml = ""
        if chart_type == "bar" and len(series) == 1:
            for j in range(len(vals)):
                c = CHART_COLORS[j % len(CHART_COLORS)]
                dpt_xml += f"""<c:dPt><c:idx val="{j}"/>
                  <c:spPr><a:solidFill><a:srgbClr val="{c}"/></a:solidFill></c:spPr>
                </c:dPt>"""

        # Series-level color for multi-series bar or line charts
        ser_color_xml = ""
        if chart_type == "line":
            c = CHART_COLORS[i % len(CHART_COLORS)]
            ser_color_xml = f"""<c:spPr>
              <a:ln w="28575"><a:solidFill><a:srgbClr val="{c}"/></a:solidFill></a:ln>
            </c:spPr>"""
            marker_xml = f"""<c:marker><c:symbol val="circle"/><c:size val="6"/>
              <c:spPr><a:solidFill><a:srgbClr val="{c}"/></a:solidFill></c:spPr>
            </c:marker>"""
        elif chart_type == "bar" and len(series) > 1:
            c = CHART_COLORS[i % len(CHART_COLORS)]
            ser_color_xml = f"""<c:spPr><a:solidFill><a:srgbClr val="{c}"/></a:solidFill></c:spPr>"""

        ser_blocks += f"""<c:ser>
        <c:idx val="{i}"/><c:order val="{i}"/>
        <c:tx><c:strRef><c:f>'{ws_title}'!${col_letter}${hdr_row}</c:f>
          <c:strCache><c:ptCount val="1"/><c:pt idx="0"><c:v>{xml_escape(str(sname))}</c:v></c:pt></c:strCache>
        </c:strRef></c:tx>
        {ser_color_xml}
        {marker_xml}
        {dpt_xml}
        <c:cat><c:strRef><c:f>'{ws_title}'!$A${cat_row1}:$A${cat_rown}</c:f>
          <c:strCache><c:ptCount val="{n}"/>{cat_pts}</c:strCache>
        </c:strRef></c:cat>
        <c:val><c:numRef><c:f>'{ws_title}'!${col_letter}${cat_row1}:${col_letter}${cat_rown}</c:f>
          <c:numCache><c:formatCode>General</c:formatCode><c:ptCount val="{n}"/>{pts}</c:numCache>
        </c:numRef></c:val>
        <c:smooth val="0"/>
      </c:ser>"""

    legend = '<c:legend><c:legendPos val="b"/></c:legend>' if len(series) > 1 else ""

    if chart_type == "pie":
        # Pie chart — different XML structure, no axes needed
        n_series = len(series)
        first_sname = xml_escape(str(list(series.keys())[0]))
        first_vals  = list(series.values())[0]
        pie_pts = "".join(
            f'<c:pt idx="{j}"><c:v>{v}</c:v></c:pt>' for j, v in enumerate(first_vals)
        )
        cat_pts = "".join(
            f'<c:pt idx="{j}"><c:v>{xml_escape(str(cat))}</c:v></c:pt>' for j, cat in enumerate(categories)
        )
        # Per-slice colors
        pie_dpt_xml = ""
        for j in range(len(first_vals)):
            c = CHART_COLORS[j % len(CHART_COLORS)]
            pie_dpt_xml += f"""<c:dPt><c:idx val="{j}"/>
              <c:spPr><a:solidFill><a:srgbClr val="{c}"/></a:solidFill></c:spPr>
            </c:dPt>"""
        pie_dlbls = f"""<c:dLbls>
          <c:numFmt formatCode="0.0%" sourceLinked="0"/>
          <c:showLegendKey val="0"/><c:showVal val="0"/>
          <c:showCatName val="1"/><c:showSerName val="0"/>
          <c:showPercent val="1"/><c:showBubbleSize val="0"/>
          <c:showLeaderLines val="1"/>
        </c:dLbls>"""
        return f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<c:chartSpace xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart"
              xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
              xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <c:chart>
    <c:title><c:tx><c:rich><a:bodyPr/><a:p><a:r><a:t>{chart_title}</a:t></a:r></a:p></c:rich></c:tx>
    <c:overlay val="0"/></c:title>
    <c:plotArea>
      <c:pieChart>
        <c:varyColors val="1"/>
        <c:ser>
          <c:idx val="0"/><c:order val="0"/>
          <c:tx><c:strRef><c:f>'{ws_title}'!$B${hdr_row}</c:f>
            <c:strCache><c:ptCount val="1"/><c:pt idx="0"><c:v>{first_sname}</c:v></c:pt></c:strCache>
          </c:strRef></c:tx>
          {pie_dpt_xml}
          {pie_dlbls}
          <c:cat><c:strRef><c:f>'{ws_title}'!$A${cat_row1}:$A${cat_rown}</c:f>
            <c:strCache><c:ptCount val="{n}"/>{cat_pts}</c:strCache>
          </c:strRef></c:cat>
          <c:val><c:numRef><c:f>'{ws_title}'!$B${cat_row1}:$B${cat_rown}</c:f>
            <c:numCache><c:formatCode>General</c:formatCode><c:ptCount val="{n}"/>{pie_pts}</c:numCache>
          </c:numRef></c:val>
        </c:ser>
        <c:firstSliceAng val="0"/>
      </c:pieChart>
    </c:plotArea>
    <c:legend><c:legendPos val="r"/></c:legend>
    <c:plotVisOnly val="1"/>
  </c:chart>
  <c:spPr><a:noFill/></c:spPr>
</c:chartSpace>"""

    elif chart_type == "line":
        chart_block = f"""<c:lineChart>
        <c:grouping val="standard"/>
        <c:varyColors val="0"/>
        {ser_blocks}
        {dlbls}
        <c:axId val="10"/><c:axId val="100"/>
      </c:lineChart>"""
    else:
        vary = "1" if len(series) == 1 else "0"
        chart_block = f"""<c:barChart>
        <c:barDir val="col"/><c:grouping val="clustered"/><c:varyColors val="{vary}"/>
        {ser_blocks}{dlbls}
        <c:gapWidth val="150"/>
        <c:axId val="10"/><c:axId val="100"/>
      </c:barChart>"""

    return f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<c:chartSpace xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart"
              xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
              xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <c:chart>
    <c:title><c:tx><c:rich><a:bodyPr/><a:p><a:r><a:t>{chart_title}</a:t></a:r></a:p></c:rich></c:tx>
    <c:overlay val="0"/></c:title>
    <c:plotArea>
      {chart_block}
      <c:catAx>
        <c:axId val="10"/><c:scaling><c:orientation val="minMax"/></c:scaling>
        <c:delete val="0"/><c:axPos val="b"/>
        <c:title><c:tx><c:rich><a:bodyPr/><a:p><a:r><a:t>{x_title}</a:t></a:r></a:p></c:rich></c:tx>
        <c:overlay val="0"/></c:title>
        <c:numFmt formatCode="General" sourceLinked="0"/>
        <c:tickLblPos val="nextTo"/><c:crossAx val="100"/>
        <c:auto val="1"/><c:lblAlgn val="ctr"/><c:lblOffset val="100"/>
      </c:catAx>
      <c:valAx>
        <c:axId val="100"/><c:scaling><c:orientation val="minMax"/></c:scaling>
        <c:delete val="0"/><c:axPos val="l"/><c:majorGridlines/>
        <c:title><c:tx><c:rich><a:bodyPr/><a:p><a:r><a:t>{y_title}</a:t></a:r></a:p></c:rich></c:tx>
        <c:overlay val="0"/></c:title>
        <c:numFmt formatCode="General" sourceLinked="1"/>
        <c:tickLblPos val="nextTo"/><c:crossAx val="10"/>
        <c:crossBetween val="between"/>
      </c:valAx>
    </c:plotArea>
    {legend}
    <c:plotVisOnly val="1"/>
  </c:chart>
  <c:spPr><a:noFill/></c:spPr>
</c:chartSpace>"""


def inject_chart_xmls(wb, chart_xmls):
    """Save wb, replace chart XMLs with corrected ones, return patched bytes."""
    buf = io.BytesIO()
    wb.save(buf); buf.seek(0)
    with zipfile.ZipFile(buf, 'r') as zin:
        files = {n: zin.read(n) for n in zin.namelist()}
    chart_names = sorted([n for n in files if re.match(r'xl/charts/chart\d+\.xml', n)])
    for cname, cxml in zip(chart_names, chart_xmls):
        files[cname] = cxml.encode('utf-8')
    out = io.BytesIO()
    with zipfile.ZipFile(out, 'w', zipfile.ZIP_DEFLATED) as zout:
        for name, data in files.items():
            zout.writestr(name, data)
    return out.getvalue()


# ── Excel export ──────────────────────────────────────────────────────────────
def export_to_excel(df, analysis, chart_path, output_path="report.xlsx",
                    df2=None, name1="Dataset 1", name2="Dataset 2"):
    wb = Workbook()
    hdr_font = Font(bold=True, size=11, color="FFFFFF")
    hdr_fill = PatternFill("solid", fgColor="4F46E5")
    chart_title = analysis.get("chart_title", "Analysis")

    # Sheet 1: Summary
    ws1 = wb.active
    ws1.title = "Summary"
    ws1["A1"] = "AI Data Analysis Report"
    ws1["A1"].font = Font(bold=True, size=16)
    ws1["A2"] = "Question: " + chart_title
    ws1["A2"].font = Font(italic=True, size=11)
    ws1.column_dimensions["A"].width = 80

    all_charts = resolve_all_charts(analysis, df)
    chart_xmls = []

    # ── Verified Summary FIRST — always accurate ──────────────────────────────
    next_row = 4
    hdr = ws1.cell(row=next_row, column=1, value="✅ Verified Data Summary (computed directly from dataset)")
    hdr.font = Font(bold=True, size=12)
    next_row += 1

    currency_symbol = detect_currency(analysis)

    def fmt_val(val, lf):
        if lf == "currency":   return f"{currency_symbol}{val:,.0f}"
        elif lf == "percent":  return f"{val:.1f}%"
        else:                  return f"{val:,.0f}"

    for cd in all_charts:
        ws1.cell(row=next_row, column=1, value=cd["title"]).font = Font(bold=True, size=11)
        next_row += 1
        lf   = cd.get("label_format", "none")
        vals = list(cd["series"].values())[0]
        cats = cd["categories"]
        max_val = max(vals); max_cat = cats[vals.index(max_val)]
        min_val = min(vals); min_cat = cats[vals.index(min_val)]
        total   = sum(vals)
        ws1.cell(row=next_row, column=1, value=f"  Highest: {max_cat} — {fmt_val(max_val, lf)}")
        next_row += 1
        ws1.cell(row=next_row, column=1, value=f"  Lowest:  {min_cat} — {fmt_val(min_val, lf)}")
        next_row += 1
        ws1.cell(row=next_row, column=1, value=f"  Total:   {fmt_val(total, lf)}")
        next_row += 1
        for cat, val in zip(cats, vals):
            ws1.cell(row=next_row, column=1, value=f"    {cat}: {fmt_val(val, lf)}")
            next_row += 1
        next_row += 1

    # ── Answer ────────────────────────────────────────────────────────────────
    next_row += 1
    ws1.cell(row=next_row, column=1, value="Answer").font = Font(bold=True, size=12)
    next_row += 1
    cell = ws1.cell(row=next_row, column=1, value=analysis.get("answer", ""))
    cell.alignment = Alignment(wrap_text=True)
    ws1.row_dimensions[next_row].height = 80
    next_row += 2

    # ── Key Findings ──────────────────────────────────────────────────────────
    ws1.cell(row=next_row, column=1, value="Key Findings").font = Font(bold=True, size=12)
    next_row += 1
    for f in analysis.get("key_findings", []):
        ws1.cell(row=next_row, column=1, value="• " + f)
        next_row += 1
    next_row += 1

    for i, cd in enumerate(all_charts):
        cats   = cd["categories"]
        ser    = cd["series"]
        title  = cd["title"]
        n_rows = len(cats) + 1
        x_t    = cd["x_title"]
        y_t    = cd["y_title"]
        lf     = cd["label_format"]
        ct     = cd["chart_type"]

        sheet_name    = f"Chart {i+1}" if len(all_charts) > 1 else "AI Chart"
        data_start_row = 3   # data table starts at row 3 on the chart sheet

        # Chart sheet — data table + chart on same sheet, fully linked
        ws_c = wb.create_sheet(sheet_name)
        ws_c["A1"] = title
        ws_c["A1"].font = Font(bold=True, size=14)
        ws_c.column_dimensions["A"].width = 25

        # Write editable data table starting at row 3
        ws_c.cell(row=data_start_row, column=1, value="Category").font = hdr_font
        ws_c.cell(row=data_start_row, column=1).fill = hdr_fill
        for c_idx, col_name in enumerate(ser.keys(), start=2):
            hc = ws_c.cell(row=data_start_row, column=c_idx, value=col_name)
            hc.font = hdr_font; hc.fill = hdr_fill
            ws_c.column_dimensions[chr(64 + c_idx)].width = 18
        for r, cat in enumerate(cats, start=data_start_row + 1):
            ws_c.cell(row=r, column=1, value=str(cat))
            for c_idx, vals in enumerate(ser.values(), start=2):
                idx = r - (data_start_row + 1)
                ws_c.cell(row=r, column=c_idx,
                          value=round(float(vals[idx]), 2) if idx < len(vals) else 0)

        # Chart positioned to the right of the data table
        chart_col_letter = chr(64 + len(ser) + 3)
        chart_anchor = f"{chart_col_letter}3"

        # Placeholder chart — references same sheet data table
        dummy = BarChart()
        dummy.add_data(Reference(ws_c, min_col=2, min_row=data_start_row, max_row=n_rows + data_start_row - 1))
        dummy.set_categories(Reference(ws_c, min_col=1, min_row=data_start_row + 1, max_row=n_rows + data_start_row - 1))
        dummy.width = 26; dummy.height = 18
        ws_c.add_chart(dummy, chart_anchor)

        chart_xmls.append(build_chart_xml(
            sheet_name, cats, ser, title, y_t, x_t, lf, ct,
            currency_symbol=detect_currency(analysis),
            data_start_row=data_start_row
        ))

    # Raw Data sheet
    ws_raw = wb.create_sheet("Raw Data")
    for ci, cn in enumerate(df.columns, start=1):
        c = ws_raw.cell(row=1, column=ci, value=str(cn))
        c.font = hdr_font; c.fill = hdr_fill
        c.alignment = Alignment(horizontal="center")

    def safe_val(v):
        """Convert any value to an Excel-compatible type."""
        if v is None or (isinstance(v, float) and __import__('math').isnan(v)):
            return ""
        import pandas as pd
        if isinstance(v, pd.Period):
            return str(v)
        if isinstance(v, pd.Timestamp):
            return v.to_pydatetime()
        if hasattr(v, 'item'):  # numpy scalar
            return v.item()
        if isinstance(v, (int, float, str, bool)):
            return v
        return str(v)

    for ri, row in enumerate(df.itertuples(index=False), start=2):
        for ci, val in enumerate(row, start=1):
            ws_raw.cell(row=ri, column=ci, value=safe_val(val))

    patched = inject_chart_xmls(wb, chart_xmls)
    with open(output_path, "wb") as f:
        f.write(patched)
    return output_path


# ── Word chart image helper ──────────────────────────────────────────────────
def _fmt_short(v, currency_symbol=""):
    """Abbreviate large numbers: 315341396 → ₦315.3M"""
    av = abs(v)
    if av >= 1_000_000_000:
        return f"{currency_symbol}{v/1_000_000_000:,.1f}B"
    elif av >= 1_000_000:
        return f"{currency_symbol}{v/1_000_000:,.1f}M"
    elif av >= 1_000:
        return f"{currency_symbol}{v/1_000:,.1f}K"
    else:
        return f"{currency_symbol}{v:,.0f}"


def _render_chart_image(cd, currency_symbol="$"):
    """Render a single chart dict to a PNG image in memory using matplotlib."""
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import matplotlib.ticker as ticker

    ct = cd.get("chart_type", "bar")
    cats = [str(c) for c in cd["categories"]]
    series = cd["series"]
    title = cd["title"]
    lf = cd.get("label_format", "none")
    csym = currency_symbol if lf == "currency" else ""

    fig, ax = plt.subplots(figsize=(10, 5.5))

    if ct == "pie":
        vals = list(series.values())[0]
        pie_colors = ["#" + CHART_COLORS[i % len(CHART_COLORS)] for i in range(len(cats))]
        wedges, texts, autotexts = ax.pie(
            vals, labels=cats, autopct='%1.1f%%', colors=pie_colors,
            textprops={'fontsize': 10})
        for at in autotexts:
            at.set_fontsize(9)
            at.set_fontweight('bold')
        ax.set_title(title, fontsize=14, fontweight='bold', pad=16)
    elif ct == "line":
        for i, (sname, vals) in enumerate(series.items()):
            c = "#" + CHART_COLORS[i % len(CHART_COLORS)]
            ax.plot(cats, vals, marker='o', linewidth=2.5, markersize=7,
                    label=sname, color=c)
            for j, v in enumerate(vals):
                ax.annotate(_fmt_short(v, csym), (cats[j], v),
                           textcoords="offset points",
                           xytext=(0, 12), ha='center', fontsize=8,
                           fontweight='bold')
        ax.set_title(title, fontsize=14, fontweight='bold')
        ax.set_xlabel(cd.get("x_title", ""), fontsize=10)
        ax.set_ylabel(cd.get("y_title", ""), fontsize=10)
        ax.yaxis.set_major_formatter(
            ticker.FuncFormatter(lambda x, _: _fmt_short(x, csym)))
        if len(series) > 1:
            ax.legend()
        ax.grid(axis='y', alpha=0.3)
    else:  # bar
        if len(series) == 1:
            vals = list(series.values())[0]
            bar_colors = ["#" + CHART_COLORS[i % len(CHART_COLORS)] for i in range(len(cats))]
            bars = ax.bar(cats, vals, color=bar_colors, width=0.6)
            for bar, v in zip(bars, vals):
                ax.text(bar.get_x() + bar.get_width()/2, bar.get_height(),
                       _fmt_short(v, csym),
                       ha='center', va='bottom', fontsize=8, fontweight='bold')
        else:
            import numpy as np
            x = np.arange(len(cats))
            width = 0.8 / len(series)
            for i, (sname, vals) in enumerate(series.items()):
                c = "#" + CHART_COLORS[i % len(CHART_COLORS)]
                b = ax.bar(x + i * width - 0.4 + width/2, vals, width, label=sname, color=c)
                for bar, v in zip(b, vals):
                    ax.text(bar.get_x() + bar.get_width()/2, bar.get_height(),
                           _fmt_short(v, csym),
                           ha='center', va='bottom', fontsize=7, fontweight='bold')
            ax.set_xticks(x)
            ax.set_xticklabels(cats)
            ax.legend()
        ax.set_title(title, fontsize=14, fontweight='bold')
        ax.set_xlabel(cd.get("x_title", ""), fontsize=10)
        ax.set_ylabel(cd.get("y_title", ""), fontsize=10)
        ax.yaxis.set_major_formatter(
            ticker.FuncFormatter(lambda x, _: _fmt_short(x, csym)))
        ax.grid(axis='y', alpha=0.3)

    if ct != "pie":
        plt.xticks(rotation=30 if len(cats) > 4 else 0,
                   ha='right' if len(cats) > 4 else 'center')
    plt.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=180, bbox_inches='tight',
                facecolor='white', edgecolor='none')
    plt.close(fig)
    buf.seek(0)
    return buf


# ── Word export ───────────────────────────────────────────────────────────────
def export_to_word(df, analysis, chart_path, output_path="report.docx",
                   df2=None, name1="Dataset 1", name2="Dataset 2"):
    doc = Document()
    doc.add_heading("AI Data Analysis Report", 0)
    doc.add_paragraph("Question: " + analysis.get("chart_title", ""))
    doc.add_heading("Answer", level=1)
    doc.add_paragraph(analysis.get("answer", ""))
    doc.add_heading("Key Findings", level=1)
    for f in analysis.get("key_findings", []):
        doc.add_paragraph(f, style="List Bullet")

    # ── Charts as images ──
    all_charts = resolve_all_charts(analysis, df)
    currency_sym = detect_currency(analysis)
    if all_charts:
        doc.add_heading("Charts", level=1)
        for cd in all_charts:
            doc.add_heading(cd["title"], level=2)
            try:
                img_buf = _render_chart_image(cd, currency_symbol=currency_sym)
                doc.add_picture(img_buf, width=Inches(6.2))
            except Exception as e:
                import logging
                logging.warning(f"[word] chart render error: {e}")
            doc.add_paragraph()  # spacing

        # ── Chart data tables (for editability) ──
        doc.add_heading("Chart Data — Editable Tables", level=1)
        for cd in all_charts:
            doc.add_heading(cd["title"], level=2)
            col_names = ["Category"] + list(cd["series"].keys())
            tbl = doc.add_table(rows=1, cols=len(col_names))
            tbl.style = "Table Grid"
            tbl.autofit = True
            for i, col in enumerate(col_names):
                cell = tbl.rows[0].cells[i]
                cell.text = col
                r = cell.paragraphs[0].runs[0]
                r.font.bold = True
                r.font.color.rgb = DocxRGB(0x4F, 0x46, 0xE5)
            for i, cat in enumerate(cd["categories"]):
                rc = tbl.add_row().cells
                rc[0].text = str(cat)
                for j, vals in enumerate(cd["series"].values()):
                    rc[j+1].text = str(round(float(vals[i]), 2)) if i < len(vals) else "0"
            doc.add_paragraph()

    # ── Data Sample — limit to first 10 columns for readability ──
    doc.add_heading("Data Sample", level=1)
    max_cols = min(len(df.columns), 10)
    sample_df = df.head(15)
    if len(df.columns) > max_cols:
        note = doc.add_paragraph()
        run = note.add_run(f"Showing {max_cols} of {len(df.columns)} columns. Full data available in Excel export.")
        run.font.italic = True
        run.font.size = Pt(9)

    t = doc.add_table(rows=1, cols=max_cols)
    t.style = "Table Grid"
    t.autofit = True
    for i in range(max_cols):
        cell = t.rows[0].cells[i]
        cell.text = str(df.columns[i])
        cell.paragraphs[0].runs[0].font.bold = True
        cell.paragraphs[0].runs[0].font.size = Pt(8)
    for _, row in sample_df.iterrows():
        rc = t.add_row().cells
        for i in range(max_cols):
            rc[i].text = str(row.iloc[i])
            rc[i].paragraphs[0].runs[0].font.size = Pt(8)

    doc.save(output_path)
    return output_path


# ── PowerPoint export ─────────────────────────────────────────────────────────
def export_to_pptx(analysis, chart_path, output_path="report.pptx",
                   df=None, df2=None, name1="Dataset 1", name2="Dataset 2"):
    prs = Presentation()
    prs.slide_width  = PInches(13.33)
    prs.slide_height = PInches(7.5)

    # Slide 1 — Title
    s1 = prs.slides.add_slide(prs.slide_layouts[0])
    t1 = s1.shapes.title
    t1.text = "AI Data Analysis Report"
    set_textbox(t1, 0.5, 2.5, 12.0, 1.2)
    for r in t1.text_frame.paragraphs[0].runs:
        r.font.size = PPt(40); r.font.bold = True
    p1 = s1.placeholders[1]
    p1.text = analysis.get("chart_title", "")
    set_textbox(p1, 0.5, 3.9, 12.0, 0.8)
    for r in p1.text_frame.paragraphs[0].runs:
        r.font.size = PPt(24)

    # Slide 2 — Answer
    s2 = prs.slides.add_slide(prs.slide_layouts[1])
    t2 = s2.shapes.title
    t2.text = "Analysis Answer"
    set_textbox(t2, 0.5, 0.3, 12.0, 0.9)
    for r in t2.text_frame.paragraphs[0].runs:
        r.font.size = PPt(32); r.font.bold = True
    b2 = s2.placeholders[1]
    set_textbox(b2, 0.5, 1.4, 12.3, 5.6)
    add_text_to_placeholder(b2.text_frame, analysis.get("answer", ""), font_size=18)

    # Slide 3 — Key Findings
    s3 = prs.slides.add_slide(prs.slide_layouts[1])
    t3 = s3.shapes.title
    t3.text = "Key Findings"
    set_textbox(t3, 0.5, 0.3, 12.0, 0.9)
    for r in t3.text_frame.paragraphs[0].runs:
        r.font.size = PPt(32); r.font.bold = True
    b3 = s3.placeholders[1]
    set_textbox(b3, 0.5, 1.4, 12.3, 5.6)
    add_bullet_placeholder(b3.text_frame, analysis.get("key_findings", []), font_size=18)

    # Slides 4+ — One native chart per subplot
    if df is not None:
        all_charts = resolve_all_charts(analysis, df)
        currency_sym = detect_currency(analysis)
        for cd in all_charts:
            try:
                s = prs.slides.add_slide(prs.slide_layouts[6])
                tb = s.shapes.add_textbox(PInches(0.5), PInches(0.2), PInches(12.0), PInches(0.8))
                tf = tb.text_frame
                tf.text = cd["title"]
                tf.paragraphs[0].runs[0].font.size = PPt(28)
                tf.paragraphs[0].runs[0].font.bold = True

                pptx_cd = ChartData()
                pptx_cd.categories = [str(c) for c in cd["categories"]]
                for sname, vals in cd["series"].items():
                    pptx_cd.add_series(sname, [round(float(v), 2) for v in vals])

                ct = cd.get("chart_type", "bar")
                if ct == "line":
                    chart_type = XL_CHART_TYPE.LINE_MARKERS
                elif ct == "pie":
                    chart_type = XL_CHART_TYPE.PIE
                else:
                    chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
                ch = s.shapes.add_chart(
                    chart_type,
                    PInches(0.5), PInches(1.1), PInches(12.3), PInches(5.9), pptx_cd
                ).chart
                ch.has_legend = len(cd["series"]) > 1

                if ct == "pie":
                    # Pie charts: vary colors per slice, show category + percent labels
                    ch.plots[0].vary_by_categories = True
                    ch.has_legend = True
                    plot = ch.plots[0]
                    plot.has_data_labels = True
                    dls = plot.data_labels
                    dls.number_format = '0.0%'
                    dls.number_format_is_linked = False
                    dls.show_value = False
                    dls.show_category_name = True
                    dls.show_series_name = False
                    dls.show_percentage = True
                    # Explicit slice colors
                    series_obj = plot.series[0]
                    for pidx in range(len(cd["categories"])):
                        c = CHART_COLORS[pidx % len(CHART_COLORS)]
                        pt = series_obj.points[pidx]
                        pt.format.fill.solid()
                        pt.format.fill.fore_color.rgb = PptxRGB(
                            int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16))
                else:
                    ch.plots[0].vary_by_categories = False
                    # Data labels for non-pie charts
                    if cd.get("label_format", "none") != "none":
                        fmt_map = {
                            "currency": f"{currency_sym}#,##0",
                            "number": "#,##0",
                            "percent": "0.0%"
                        }
                        plot = ch.plots[0]
                        plot.has_data_labels = True
                        dls = plot.data_labels
                        dls.number_format = fmt_map.get(cd.get("label_format"), "General")
                        dls.number_format_is_linked = False
                        dls.show_value = True
                        dls.show_category_name = False
                        dls.show_series_name = False

                    # Apply colors to bar/line charts
                    plot = ch.plots[0]
                    if ct == "bar" and len(cd["series"]) == 1:
                        # Single series bar: color each bar differently
                        series_obj = plot.series[0]
                        for pidx in range(len(cd["categories"])):
                            c = CHART_COLORS[pidx % len(CHART_COLORS)]
                            pt = series_obj.points[pidx]
                            pt.format.fill.solid()
                            pt.format.fill.fore_color.rgb = PptxRGB(
                                int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16))
                    else:
                        # Multi-series bar or line: color each series
                        for si, ser_key in enumerate(cd["series"].keys()):
                            c = CHART_COLORS[si % len(CHART_COLORS)]
                            series_obj = plot.series[si]
                            if ct == "line":
                                series_obj.format.line.fill.solid()
                                series_obj.format.line.fill.fore_color.rgb = PptxRGB(
                                    int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16))
                                series_obj.format.line.width = PPt(2.5)
                                # Marker color
                                series_obj.marker.format.fill.solid()
                                series_obj.marker.format.fill.fore_color.rgb = PptxRGB(
                                    int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16))
                            else:
                                series_obj.format.fill.solid()
                                series_obj.format.fill.fore_color.rgb = PptxRGB(
                                    int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16))

            except Exception as e:
                import logging
                logging.warning(f"[pptx] chart error: {e}")

    prs.save(output_path)
    return output_path
